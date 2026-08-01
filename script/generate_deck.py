import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette Constants
    COLOR_BG = RGBColor(0x0B, 0x10, 0x20)
    COLOR_CARD = RGBColor(0x11, 0x18, 0x27)
    COLOR_BLUE = RGBColor(0x25, 0x63, 0xEB)
    COLOR_PURPLE = RGBColor(0x7C, 0x3A, 0xED)
    COLOR_CYAN = RGBColor(0x06, 0xB6, 0xD4)
    COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    COLOR_LIGHT = RGBColor(0xD1, 0xD5, 0xDB)
    COLOR_MUTED = RGBColor(0x9C, 0xA3, 0xAF)
    COLOR_SUCCESS = RGBColor(0x10, 0xB9, 0x81)
    COLOR_WARNING = RGBColor(0xF5, 0x9E, 0x0B)
    COLOR_DANGER = RGBColor(0xEF, 0x44, 0x44)

    FONT_NAME = 'Vazirmatn'

    slides_data = [
        {
            "type": "hero",
            "tag": "کۆرسی ژیری دەستکرد • ٢ کاتژمێر",
            "title": "ژیری دەستکرد",
            "subtitle": "«ژیری دەستکرد لە ژیانی ڕۆژانە و کاردا» — ڕێنمایی گشتگیر بۆ قوتابی، مامۆستا، کارمەند، بزنس و گەشەپێدەران",
            "presenter": "پێشکەشکار: ئەندازیار عبدالرحمن اسماعیل",
            "notes": "بەخێربێن بۆ کۆرسی ژیری دەستکرد. ئەمڕۆ تیشک دەخەینە سەر بەکارهێنانی کردارییانەی AI لە ژیانی ڕۆژانە و سەرجەم سێکتەرەکانی کاردا."
        },
        {
            "type": "presenter_intro",
            "category": "پێشەکی",
            "title": "دەربارەی پێشکەشکار (خۆناساندن)",
            "name": "ئەندازیار عبدالرحمن اسماعیل",
            "roles": "IT Manager  •  Web Developer  •  Technical Trainer",
            "image_path": r"c:\Users\kurdn\Desktop\Semenar ChatGPT\assets\image\agha.png",
            "education": [
                "ئەندازیاری سیستەمی زانیاری",
                "تەکنەلۆجیای زانیاری (IT)"
            ],
            "experience": [
                "زیاتر لە ٥ ساڵ ئەزموونی وانەوتنەوە و ڕاهێنەرایەتی لە چەندین بواری جیاجیادا.",
                "بەڕێوەبەری پێشووی بەشی IT لە چەندین کۆمپانیا.",
                "ئەزموون لە بواری جیاوازی سەر بە IT."
            ],
            "notes": "دەربارەی پێشکەشکاری سیمینارەکە، ئەندازیار عبدالرحمن اسماعیل، پیشە، بڕوانامەکان و ئەزموونی کاری لە بواری تەکنەلۆجیادا."
        },
        {
            "type": "grid6_agenda",
            "category": "نەخشەی رێگای کۆرسەکە",
            "title": "بابەتە سەرەکیەکانی کۆرسەکە (١٢ بەش)",
            "items": [
                {"num": "بەشی ١ - ٣", "title": "بنەما, مۆدێلەکان & پرۆمپت", "body": "ژیری دەستکرد چییە؟ جۆرەکانی مۆدێلەکان، و ئەندازیاری پرۆمپتی پیشەیی."},
                {"num": "بەشی ٤ - ٦", "title": "قوتابی، مامۆستا & خاوەنکار", "body": "کورتکردنەوەی PDF، داڕشتنی پلان، فیدباک و ڕاپۆرتی ئۆتۆماتیک."},
                {"num": "بەشی ٧ - ٨", "title": "بزنس & گەشەپێدەر", "body": "پرۆسەی کاری ١١ هەنگاوی بزنس و گەشەپێدانی کۆد بە AI."},
                {"num": "☕ پشوو", "title": "کۆفی برێک (١٥ خولەک)", "body": "پشوو پاش تەواوبوونی وانەکانی خۆکارکردن، بزنس و گەشەپێدان."},
                {"num": "بەشی ٩ - ١٠", "title": "دیزاینەر & سنوورەکانی AI", "body": "دروستکردنی وێنە و ڤیدیۆ، و تێگەیشتن لە هەڵەکان و دیپ فەیک."},
                {"num": "بەشی ١١ - ١٢", "title": "نموونەی ڕاستەوخۆ & داهاتوو", "body": "تەماشاکردنی دێمۆی ڕاستەوخۆ و داڕشتنی پاشەرۆژی ئەی ئای."}
            ],
            "notes": "خشتەی کۆرسەکە بە شێوازێک داڕێژراوە کە لە سادەترین بنەماوە دەستپێبکات تا عەمەلیترین دیزاین و کۆد."
        },
        {
            "type": "grid5",
            "category": "بەشی ١ (١٠ - ١٥ خولەک)",
            "title": "ژیری دەستکرد چییە؟ (تێگەیشتن لە بنەماکان)",
            "items": [
                {"title": "ژیری دەستکرد چییە؟", "body": "سیستەمێک کە بیرکاری و داتا بەکاردێنێت بۆ شیکردنەوە و بڕیاردان وەک مرۆڤ."},
                {"title": "لە کوێ دەبینین؟", "body": "لە تەلیفۆن، گووگڵ، فەیسبووک، سیستەمی پزیشکی و ئۆتۆمبێلی خۆڕەودا."},
                {"title": "چۆن بیر دەکات؟", "body": "بە شیکردنەوەی ملیارەها نەخش و داتا لە ڕێگەی تۆڕە دەمارییەکانەوە (Neural Networks)."},
                {"title": "مۆدێلی زمانی چییە؟", "body": "مۆدێلی زمانی گەورە (Large Language Model) کە لە دەق تێدەگات و نووسین دەخولقێنێت."},
                {"title": "چۆن فێر دەبێت؟", "body": "لە ڕێگەی ڕاهێنان (Training) لەسەر تێکست و زانیاری ڕابردوو؛ «جادوو نییە، بیرکارییە!»"}
            ],
            "notes": "ئامانجی سەرەکی ئەم بەشە ئەوەیە بەشداربووان تێبگەن AI جادوو نییە، بەڵکو تەکنەلۆجیایەکی شیکارییە."
        },
        {
            "type": "grid6_agenda",
            "category": "بەشی ٢ (١٥ خولەک)",
            "title": "جۆرەکانی ئامرازی ژیری دەستکرد و بەکارهێنانیان",
            "items": [
                {"num": "چات جیپیتی (ChatGPT)", "title": "نووسین، کۆد و پلان", "body": "باشترین بۆ: نووسینی دەق، شیکاری، کۆدکردن، و داڕشتنی پلانی گشتی."},
                {"num": "کلۆد (Claude)", "title": "PDF، ڕاپۆرت و نووسین", "body": "باشترین بۆ: خوێندنەوەی فایلی PDF ی درێژ، نووسینی فەرمی و ڕاپۆرتی زانستی."},
                {"num": "جێمینای (Gemini)", "title": "گووگڵ، Docs & Gmail", "body": "باشترین بۆ: بەستنەوە بە بەڵگەنامەکانی گووگڵ، گەڕان و خزمەتگوزارییەکان."},
                {"num": "پێرپلێکسیتی (Perplexity)", "title": "توێژینەوە و سەرچاوەکان", "body": "باشترین بۆ: لێکۆڵینەوەی ئەکادیمی، بەڵگە و هێنانی لینک و سەرچاوەی ڕاستەقینە."},
                {"num": "کۆپایڵۆت (Copilot)", "title": "ماڵپەڕەکانی مایکرۆسۆفت", "body": "باشترین بۆ: ئۆتۆماتیکردنی Excel، Word، و دروستکردنی پرێزێنتەیشن لە PowerPoint."},
                {"num": "گرۆک (Grok)", "title": "تۆڕی ئێکس (Twitter)", "body": "باشترین بۆ: زانیارییەکانی ڕووداوە زۆر تازەکان لە سۆشیال میدیاوە."}
            ],
            "notes": "بۆ هەر ئەرکێک ئامرازێکی گونجاو هەیە؛ بەکارهێنانی ئامرازی ڕاست خێرایی کارەکان دەقات دەکاتەوە."
        },
        {
            "type": "comparison",
            "category": "بەشی ٣ (٢٥ خولەک)",
            "title": "ئەندازیاری پرۆمپت — ئاستەکانی نووسینی پرۆمپت",
            "weak_title": "پرۆمپتی لاواز",
            "weak_code": "«بابەتێکم بۆ بنووسە.»",
            "weak_body": "گشتییە و دەرئەنجامێکی لاواز دەدات.",
            "strong_title": "پرۆمپتی زۆر باش (پێشکەوتوو)",
            "strong_code": "«تۆ مامۆستای ژیری دەستکردیت. بابەتێکی ٥٠٠ وشەی فەرمی بنووسە بۆ قوتابی زانکۆ، بە زمانی کوردی، بە سەرناو، نموونە، و کۆتایی، و ٥ پرسیاری تاقیکردنەوە لە کۆتایی زیاد بکە.»",
            "strong_body": "ڕۆڵ، ئۆدیانس، فۆرمات، زمانی دیاریکراو و ئەرکی تەواو دەدات بە AI.",
            "notes": "ئەم بەشە گرنگترین بەشی کۆرسەکەیە! فێربوونی دروستکردنی پرۆمپت کۆنتڕۆڵی تەواوت دەداتێ."
        },
        {
            "type": "grid5",
            "category": "نموونەی پرۆمپت",
            "title": "بانکی نموونەکان: فێرکاری، خوێندن، بزنس و کۆد",
            "items": [
                {"title": "١. بواری فێرکاری (مامۆستا)", "body": "🔴 لاواز: «تاقیکردنەوەیەکم بۆ بکە.»\n🟡 ناوەند: «تاقیکردنەوەیەک بۆ پۆلی ٩ بە ٥ پرسیار.»\n🟢 باش: «تۆ مامۆستایەکی بیرکاری لێهاتووی پۆلی نۆیت. تاقیکردنەوەیەکی فرەبژاردە لەسەر جەبر دروستبکە...»"},
                {"title": "٢. بواری خوێندن (قوتابی)", "body": "🔴 لاواز: «بابەتەکەم بۆ ڕوونبکەرەوە.»\n🟡 ناوەند: «ژیری دەستکردم بۆ کورت بکەرەوە.»\n🟢 باش: «تۆ خوێندکاری زانکۆیت لە بەشی کۆمپیوتەر. چەمکی Neural Networks کورت بکەرەوە بە...»"},
                {"title": "٣. بواری کار (بزنس)", "body": "🔴 لاواز: «ڕیکلامێکم بۆ بنووسە.»\n🟡 ناوەند: «پۆستێکی فەیسبووک بۆ ئایفۆن.»\n🟢 باش: «تۆ پسپۆڕی مارکێتینگیت. پۆستێکی سەرنجڕاکێش بنووسە بۆ فەیسبووک بۆ فرۆشتنی مۆبایل...»"},
                {"title": "٤. بواری پڕۆگرامسازی", "body": "🔴 لاواز: «کۆدی HTMLم دەوێت.»\n🟡 ناوەند: «دوگمەیەکی جوان دروستبکە.»\n🟢 باش: «تۆ گەشەپێدەرێکی لێهاتووی وێبیت. کۆدێکی پاکی HTML و CSS دروستبکە بۆ دوگمەیەکی...»"}
            ],
            "notes": "بانکی نموونەکانی پرۆمپت پیشانی دەدات چۆن پرۆمپتەکان لە لاوازەوە دەگۆڕدرێن بۆ ناوەند و پاشان بۆ زۆر باش."
        },
        {
            "type": "grid5",
            "category": "داڕشتنی پرۆمپت",
            "title": "پێکهاتەی زێڕینی پرۆمپت (ڕۆڵ، ئەرک، پاشبنەما، فۆرمات، تۆنی دەنگ)",
            "items": [
                {"title": "١. ڕۆڵ (Role)", "body": "«وەک شارەزایەک بجوڵێوە»\nئەزموون و لێهاتوویی مۆدێلەکە دیاری دەکات بۆ ئەوەی وەڵامەکە بەپێی پسپۆڕییەکە شیکاری بکات."},
                {"title": "٢. ئەرک (Task)", "body": "«دەمەوێت ئەمەم بۆ بکەیت»\nئامانجی کارەکە بە وردی دەنووسێت (وەک نووسینی نامە، شیتاڵکردنی داتا یان کورتکردنەوە)."},
                {"title": "٣. پاشبنەما (Context)", "body": "«ئەم زانیارییانەت پێویستە»\nکێ ئۆدیانسە، بۆچی دەتەوێت و هەر زانیاری پاشبنەمای پێشوو."},
                {"title": "٤. فۆرمات (Format)", "body": "«شێوازی پیشاندانی ئەنجام»\nوەک خشتە (Table)، لیستی خاڵبەندی (Bullet points)، یان پاراگرافی کورت."},
                {"title": "٥. تۆن & سنووردارکردن", "body": "«شێوازی دەربڕین و سنوور»\nتۆنی فەرمی یان سادە، لەگەڵ سنووری ژمارەی وشە و زمانی کوردی سۆرانی."}
            ],
            "notes": "داڕشتنی پرۆمپت بەم ٥ بەشە دڵنیایی سەدا سەد دەدات لە وەرگرتنی وەڵامی دروست."
        },
        {
            "type": "workflow_diag",
            "category": "بەشی ٤",
            "title": "ژیری دەستکرد بۆ قوتابی (پرۆسەی کاری قوتابی)",
            "nodes": [
                {"step": "١", "title": "فایلی PDF / دەق", "body": "داخڵکردنی فایلی وانەکە"},
                {"step": "٢", "title": "کورتکردنەوە", "body": "پوختەی وانەکە بە خاڵ"},
                {"step": "٣", "title": "نەخشەی مێشک & پرسیار", "body": "Mind Map و تاقیکردنەوە"},
                {"step": "٤", "title": "ئامادەکاری تاقیکردنەوە", "body": "فلاش کارد و پێداچوونەوە"}
            ],
            "notes": "قوتابیان دەتوانن کاتی خوێندنەوەکەیان نیوە بکەنەوە و باشتر تێبگەن بەم workflow ە."
        },
        {
            "type": "grid5",
            "category": "نموونەی پرۆمپت بۆ قوتابیان",
            "title": "بانکی پرۆمپت بۆ قوتابیان (پۆلی ١٢، زانکۆ و فێربوونی زمان)",
            "items": [
                {"title": "١. پۆلی ١٢ی ئامادەیی", "body": "«تۆ مامۆستایەکی بیرکاری لێهاتووی پۆلی ١٢ی ئامادەییت. بەندی سێیەم (داتاشراو) لە بابەتی بیرکاری بە کوردی کورت بکەرەوە. پێناسە سەرەکییەکان و سێ نموونەی ڕوونم بۆ بنووسە بە شێوازی هەنگاو بە هەنگاو کە تێگەیشتنی ئاسان بێت.»"},
                {"title": "٢. خوێندکاری زانکۆ", "body": "«تۆ خوێندکاری زانکۆیت لە بەشی زانستی کۆمپیوتەر. پوختەیەکی چڕوپڕ و زانستی بۆ توێژینەوەی 'Deep Learning' بە زمانی کوردی بنووسە کە پێکبێت لە کورتە، مێتۆدۆلۆژی، و گرنگترین ئەنجامەکان بە خاڵ بۆ سڵایدی پرێزێنتەیشن.»"},
                {"title": "٣. فێربوونی زمان", "body": "«تۆ مامۆستایەکی زمانی ئینگلیزیت بۆ ئاستی مامناوەند (B1). لیستێک لە ١٠ دەربڕین و فۆرمۆڵی گفتوگۆ لەسەر گەشتکردن و فڕینی فڕۆکە بنووسە بە ئینگلیزی، لەگەڵ مانا و ڕوونکردنەوەی کوردی، و گفتوگۆیەکی نموونەیی.»"}
            ],
            "notes": "نموونەی پرۆمپتە ئاسایی و گونجاوەکان بۆ جۆرە جیاوازەکانی خوێندکاران و قوتابیان بۆ بەکارهێنانی کرداری بە زمانی کوردی."
        },
        {
            "type": "workflow_diag",
            "category": "بەشی ٥",
            "title": "ژیری دەستکرد بۆ مامۆستا (پرۆسەی کاری مامۆستا)",
            "nodes": [
                {"step": "١", "title": "پلانی وانە", "body": "داڕشتنی پلانی وانە (Lesson Plan)"},
                {"step": "٢", "title": "سڵایدەکانی پرێزێنتەیشن", "body": "دروستکردنی سڵایدەکان"},
                {"step": "٣", "title": "پرسیار & پێوەرەکان", "body": "پرسیار و پێوەری هەڵسەنگاندن"},
                {"step": "٤", "title": "وەڵامی نموونەیی", "body": "وەڵامی ڕاستی تاقیکردنەوە"}
            ],
            "notes": "مامۆستایان دەتوانن پرۆسەی ئامادەکردنی وانە و تاقیکردنەوەکان زۆر ئاسانتر بکەن."
        },
        {
            "type": "grid5",
            "category": "نموونەی پرۆمپت بۆ مامۆستایان",
            "title": "بانکی پرۆمپت بۆ مامۆستایان (پلان، ڕێگەکانی وانەوتنەوە و هەڵسەنگاندن)",
            "items": [
                {"title": "١. دروستکردنی پلان & وانە", "body": "«تۆ شارەزایەکی داڕشتنی مەنهەج و وانەوتناوەیت لە وەزارەتی پەروەردەی هەرێمی کوردستان. پلانێکی وانەی تەواو بۆ بابەتی [ناوی بابەت بنووسە] بنووسە بۆ یەک وانەی ٤٥ خولەکی بە هەموو بەشەکانیەوە...»"},
                {"title": "٢. ڕێگەکانی وانەوتنەوە", "body": "«تۆ ڕاهێنەرێکی گەشەپێدانی مامۆستایانیت لە کوردستان. پێشنیاری ٣ ڕێگای مۆدێرن و کارلێکەر بکە بۆ شەرحکردنی بابەتی [بابەتەکە] بە شێوازێک کە قوتابی سەرنجی لای وانەکە بێت...»"},
                {"title": "٣. هەڵسەنگاندن & پرسیار", "body": "«تۆ پسپۆڕی ئامادەکردنی پرسیاری وزاریت لە هەرێم. پێنج پرسیاری فیکری لەسەر بەندی [ناوی بەند] دروستبکە و بەپێی ئاست دابەشیان بکە (ئاسان، مامناوەند، گران) لەگەڵ وەڵامە نموونەییەکان...»"}
            ],
            "notes": "نموونەی پرۆمپتە گرنگ و وزارییەکان بۆ مامۆستایانی هەرێمی کوردستان بۆ بەکارهێنانی کارایانەی AI لە ئامادەکردنی وانەکاندا."
        },
        {
            "type": "workflow_diag",
            "category": "بەشی ٦",
            "title": "ژیری دەستکرد بۆ خاوەنکار (پرۆسەی کاری خاوەنکار)",
            "nodes": [
                {"step": "١", "title": "داڕشتنی ئەرکەکان", "body": "دابەشکردن و ڕوونکردنەوەی ئەرکەکان بە AI"},
                {"step": "٢", "title": "شیتاڵکردنی ڕاپۆرتەکان", "body": "شیکردنەوەی داتای دارایی و فرۆشتن"},
                {"step": "٣", "title": "پێوەری هەڵسەنگاندن", "body": "دروستکردنی پێوەر بۆ ئەدای کارمەندان"},
                {"step": "٤", "title": "پێدانی فیدباکی بەهێز", "body": "نووسینی وەڵام و فیدباکی ڕێنمایی"}
            ],
            "notes": "خاوەنکاران دەتوانن کاتی چاودێری و بەڕێوەبردن کەم بکەنەوە و بە خێرایی بڕیاری ستراتیژی بدەن بەم workflow ە."
        },
        {
            "type": "grid5",
            "category": "نموونەی پرۆمپت بۆ خاوەنکاران",
            "title": "بانکی پرۆمپت بۆ خاوەنکاران (ئەرک، شیکردنەوەی داتا و فیدباک)",
            "items": [
                {"title": "١. داڕشتنی ئەرکەکان", "body": "«تۆ بەڕێوبەرێکی لێهاتووی تیمی کارکردنیت. دەمەوێت ئەرکی [ناوی ئەرک] ڕوون بکەیتەوە بۆ کارمەندێکی تیمەکەم بە کوردی: کورتەی ئەرک، هەنگاوەکانی ڕاپەڕاندنی بە زنجیرە، کاتی ڕادەستکردن، و پێوەری قبوڵکردن...»"},
                {"title": "٢. شیکردنەوەی داتا", "body": "«تۆ ڕاوێژکارێکی ستراتیژی بزنسیت. ئەم داتایە دارایی یان فرۆشتنەم بۆ شیکار بکە. گرنگترین ٣ خاڵی بەهێز، ٣ کێشە، و ٣ پێشنیاری کرداری کورت بۆ مانگی داهاتوو دەربکە...»"},
                {"title": "٣. فیدباکی کارمەند", "body": "«تۆ بەڕێوبەری سەرچاوە مرۆییەکان (HR)یت. فیدباکێکی فەرمی بۆ کارمەندێک بنووسە بە کوردی کە کێشەی [کێشەکە] هەیە. فیدباکەکە بە شێوازی ساندویچ بێت (ستایش، باسی کێشەکە، و هاندانی کۆتایی)...»"}
            ],
            "notes": "نموونەی پرۆمپتە کارگێڕی و پڕۆفیشناڵەکان بۆ خاوەنکاران و بەڕێوەبەران بۆ بەڕێوەبردنی باشتر و زیرەکتری تیم و بزنس."
        },
        {
            "type": "grid5",
            "category": "بەشی ٧ (بزنس)",
            "title": "ژیری دەستکرد بۆ بزنس (پرۆسەی کاری بزنس)",
            "items": [
                {"title": "بیرۆکە & پلانی بزنس", "body": "پلانی بازرگانی، وەرگرتنی ناو و لۆگۆ."},
                {"title": "مینۆ & ناسنامە", "body": "داڕشتنی کاڵا، مینۆ و براندینگ."},
                {"title": "سۆشیال میدیا & ڕیکلام", "body": "نووسینی دەقی ڕیکلام و دروستکردنی وێنە/ڤیدیۆ."},
                {"title": "وەڵامدانی کڕیار", "body": "ئۆتۆماتیکردنی وەڵامی کڕیاران و پێداچوونەوە."},
                {"title": "پلانی گەشەپێدان", "body": "شیکاری دارایی و ڕاپۆرتی گەشەپێدان."}
            ],
            "notes": "ئەم بەشە زۆر سەرنجڕاکێشە چونکە خاوەن بزنسەکان فێردەبن چۆن تەواوی مارکێتینگ و پرۆسەکانیان ئۆتۆماتیک بکەن."
        },
        {
            "type": "grid5",
            "category": "نموونەی پرۆمپت بۆ بزنس",
            "title": "بانکی پرۆمپت بۆ بزنس (پلاندانان، ڕیکلام و پشتگیری)",
            "items": [
                {"title": "١. داڕشتنی پلان & ناو", "body": "«تۆ ڕاوێژکارێکی داهێنەری بزنسیت. دەمەوێت بزنسێکی نوێی [جۆری بزنس] لە شاری [شارەکە] دابمەزرێنم. پێشنیاری ٥ ناوی کوردی، دروشمێکی بازرگانی و لیستێک لە ٣ کێشە و چارەسەریان لە بازاڕدا بکە...»"},
                {"title": "٢. نووسینی ڕیکلام", "body": "«تۆ کۆپیرایتەرێکی لێهاتوویت. دەقێکی ڕیکلامی فەیسبووک و ئینستاگرام بە کوردی بۆ فرۆشتنی [کاڵاکە] بنووسە. بە پرسیار دەستپێبکات، باسی سوودەکان بکات، CTA و ئێمۆجی تێدابێت...»"},
                {"title": "٣. پشتگیری کڕیاران", "body": "«تۆ نوێنەری پشتگیری پیشەیی کڕیاریت. نامەیەکی ڕێزدارانەی پۆزشن بە کوردی بنووسە بۆ کڕیارێک کە کێشەی [کێشەکە] هەیە. داوای لێبوردن، پێشنیارکردنی چارەسەر و دیاری قەرەبووکردنەوە تێدابێت...»"}
            ],
            "notes": "نموونەی پرۆمپتە گرنگەکانی مارکێتینگ و پەیوەندی بە کڕیارانەوە بۆ خاوەن بزنسەکان بۆ ئۆتۆماتیکردنی و بەرزکردنەوەی فرۆشتن بە کوردی."
        },
        {
            "type": "workflow_diag",
            "category": "بەشی ٨",
            "title": "ژیری دەستکرد بۆ گەشەپێدەر (پرۆسەی کاری گەشەپێدەر)",
            "nodes": [
                {"step": "١", "title": "بیرۆکە & پرۆمپت", "body": "داواکاری سیستەمەکە"},
                {"step": "٢", "title": "کۆدکردن بە AI", "body": "Flutter, Laravel, Python"},
                {"step": "٣", "title": "چاککردنی هەڵەکان", "body": "Debug & Refactor"},
                {"step": "٤", "title": "بڵاوکردنەوە & بەڵگەنامە", "body": "Deploy & Docs"}
            ],
            "notes": "گەشەپێدەران دەتوانن خێرایی کۆدنووسین و چاککردنی هەڵەکان بە ڕێژەی 3x زیاد بکەن."
        },
        {
            "type": "hero_break",
            "title": "کۆفی برێک (١٥ خولەک)",
            "subtitle": "پاش گەڕانەوە دەست دەکەین بە بەکارهێنانی AI بۆ دیزاینەر، تێگەیشتن لە مەترسییەکان و پێشاندانی دێمۆی ڕاستەوخۆ.",
            "notes": "کاتی پشوو و گۆڕینەوەی بیرۆکەکان. پاش گەڕانەوە دێینە سەر بەشە کردارییەکان."
        },
        {
            "type": "grid5",
            "category": "بەشی ٩",
            "title": "ژیری دەستکرد بۆ دیزاینەر (ئامرازەکانی دیزاین و وێنەکان)",
            "items": [
                {"title": "میدجۆرنی / چات جیپیتی", "body": "دروستکردنی بیرۆکە و وێنەی گرافیکی بەرز."},
                {"title": "ئیمەیجن / ڤیۆ (Imagen / Veo)", "body": "دروستکردنی وێنەی واقیعی و ڤیدیۆ بە AI."},
                {"title": "کانڤا (Canva AI)", "body": "دیزاینی پۆستەر و سۆشیال میدیا بە خێرایی."},
                {"title": "دەنگی دەستکرد (AI Voice)", "body": "دروستکردنی دەنگی سروشتی بۆ ڕیکلام."},
                {"title": "لۆگۆ & براندینگ", "body": "دیزاینی ڕەمزی بازرگانی و پۆستەر."}
            ],
            "notes": "ئامرازەکانی دیزاین بە AI دەرفەت دەدەن کەمپینی بێ وێنە بە کاتێکی کەم دروست بکرێت."
        },
        {
            "type": "grid5",
            "category": "نموونەی پرۆمپت بۆ دیزاینەران",
            "title": "بانکی پرۆمپت بۆ دیزاین و دروستکردنی میدیا (وێنە و ڤیدیۆی دڵ)",
            "items": [
                {"title": "١. دروستکردنی وێنە", "body": "«A highly detailed medical illustration of the human heart, showing the chambers, aorta. Scientifically accurate, 3D render, realistic textures, labeled, clean dark background --ar 16:9»"},
                {"title": "٢. گۆڕین بۆ ڤیدیۆ", "body": "«Medically accurate animation of a human heart beating rhythmically in slow motion, camera panning slowly around the chambers, highly detailed CGI, 60fps»"},
                {"title": "٣. دیزاین و پۆستەر", "body": "«A modern educational poster layout for a biology class about the heart anatomy. Minimalist design, clean typography, space for kurdish text, vector icons --ar 3:4»"}
            ],
            "notes": "نموونەی پرۆمپتەکانی دروستکردنی وێنە و ئەنیمەیشنی جوڵاوی دڵی مرۆڤ بۆ زیندەوەرزانی پۆلی ١٢ بە یارمەتی AI دیزاین کراون."
        },
        {
            "type": "grid5",
            "category": "بەشی ١٠",
            "title": "کاریگەری و سنوورەکانی AI (مەترسی و ڕەوشتەکان)",
            "items": [
                {"title": "دروستکردنی وەهم (Hallucination)", "body": "AI هەندێک جار زانیاری هەڵە دروست دەکات بە زاهیرێکی ڕاست!"},
                {"title": "مافی کۆپی (Copyright)", "body": "خاوەندارێتی داتا و وێنە دروستکراوەکان بە AI."},
                {"title": "تایبەتمەندی (Privacy)", "body": "پاراستنی داتای هەستیار و زانیاری کڕیار لە مۆدێلەکان."},
                {"title": "دیپ فەیک & هەواڵی ساختە", "body": "ڤیدیۆ و دەنگی ساختە، و ناسینەوەی زانیاری هەڵە."},
                {"title": "پشکنینی مرۆڤ", "body": "هەمیشە دەبێت مرۆڤ پێداچوونەوە بۆ دەرئەنجامەکانی AI بکات."}
            ],
            "notes": "گرنگە بەشداران بزانن AI هەمیشە ڕاست نییە و پێویستی بە پشکنین و بەرپرسیارییەتی هەیە."
        },
        {
            "type": "grid6_agenda",
            "category": "بەشی ١١",
            "title": "نموونەی ڕاستەوخۆ (Live Demos)",
            "items": [
                {"num": "کورتکردنەوەی PDF", "title": "کورتکردنەوەی فایلی درێژ", "body": "دێمۆی خوێندنەوەی PDF لە چەند چركەیەکدا."},
                {"num": "ئیمەیلی فەرمی", "title": "ئیمەیلی فەرمی", "body": "نووسینی ئیمەیل بە تۆنی گونجاو."},
                {"num": "پاوەرپۆینت", "title": "سڵایدی خۆکار", "body": "دروستکردنی پرێزێنتەیشن لە دەقەوە."},
                {"num": "پلانی بزنس", "title": "مارکێتینگ و پلانی بزنس", "body": "داڕشتنی کەمپینی فرۆشتن."},
                {"num": "کۆدی ئەپڵیکەیشن", "title": "کۆدی ئەپڵیکەیشن", "body": "نووسین و دەپلووی کۆد بە Flutter."},
                {"num": "وێنە و ڤیدیۆ", "title": "دروستکردنی میدیا", "body": "وێنە و ڤیدیۆ بە پرۆمپت."}
            ],
            "notes": "دێمۆی ڕاستەوخۆ بەشداران بە چاوی خۆیان کارایی ئامرازەکان دەبینن."
        },
        {
            "type": "grid6_agenda",
            "category": "نموونەی پرۆمپتی دێمۆکان",
            "title": "بانکی پرۆمپت بۆ دێمۆ ڕاستەوخۆکان (کارمەند، پلان، دیزاین و فێرکاری)",
            "items": [
                {"num": "١. کورتکردنەوەی PDF", "title": "کورتکردنەوەی فایلی درێژ", "body": "«تۆ شیکەرەوەی فایلی پی دی ئێفیت. ئەم فایلە بە تەواوی بخوێنەرەوە و کورت بکەرەوە بە شێوازی ٥ خاڵی سەرەکی بە کوردی...»"},
                {"num": "٢. ئیمەیلی فەرمی", "title": "نووسینی ئیمەیل", "body": "«تۆ نووسەری ئیمەیلی فەرمیت. نامەیەکی فەرمی بە کوردی بۆ دواخستنی کاتی کۆبوونەوە بنووسە لەگەڵ هۆکارەکە بە کورتی...»"},
                {"num": "٣. پاوەرپۆینت", "title": "سڵایدی خۆکار", "body": "«تۆ پسپۆڕی سڵایدەکانت. بابەتێکی کورت بکەرەوە بۆ نووسینی دەقی ٥ سڵایدی پاوەرپۆینت بە ناونیشان و دەقی کورت بە کوردی...»"},
                {"num": "٤. پلانی بزنس", "title": "پلانی کارکردن", "body": "«تۆ داڕێژەری پلانی ستراتیژی بازرگانیت. بیرۆکەی پرۆژەی [پرۆژەکە] شیتاڵ بکە بۆ دروستکردنی پلانی سەرەتایی کار بە کوردی...»"},
                {"num": "٥. کۆدی ئەپڵیکەیشن", "title": "کۆدکردنی Flutter", "body": "«تۆ گەشەپێدەرێکی لێهاتووی Flutterیت. کۆدێکی پاک و کارا دروستبکە بۆ پیشاندانی لیستێک بە ListView.builder...»"},
                {"num": "٦. وێنە & ڤیدیۆ", "title": "میدیا بە پرۆمپت", "body": "«A cinematic shot of a futuristic classroom in Kurdistan, students interacting with 3D holograms, warm lighting, Unreal Engine 5 render...»"}
            ],
            "notes": "نموونە پرۆمپتە گرنگەکان بۆ دێمۆی ڕاستەوخۆ لە کاتی پێشکەشکردنی سیمینارەکەدا بۆ وێناکردنی کارایی تەواوی ئامرازەکان."
        },
        {
            "type": "workflow_diag",
            "category": "دەرئەنجامی فێربوون",
            "title": "دەرئەنجام — ئێستا دەبێت چی فێربووبین؟ (هەنگاو بە هەنگاو)",
            "nodes": [
                {"step": "١", "title": " بیرۆکە", "body": "پێش قسەکردن لەگەڵ AI، دەبێت بیرۆکە و ئامانجی ڕوونت هەبێت."},
                {"step": "٢", "title": "پسپۆڕی", "body": "بیرکردنەوەی قوتابی 🎓، مامۆستا 👨‍🏫، بزنس 💼، ئەندازیار 💻..."},
                {"step": "٣", "title": "نووسینی پرۆمپت", "body": "دانانی Role + Task + Context + Format + Tone"},
                {"step": "٤", "title": "گفتوگۆ & دەستکاری", "body": "ناردنی پرۆمپت، پشکنینی مرۆیی و باشکردنی وەڵامەکە"}
            ],
            "notes": "ئەم سلایدە کورتکراوەی کۆرسەکەیە: پێش قسەکردن لەگەڵ AI دەبێت بیرۆکەت هەبێت، گۆشەنیگاکەت لە ئاستی پسپۆڕییەکەت دیاریبکەیت، پرۆمپتی زێڕین دابڕێژیت، پاشان پرۆسەی گفتوگۆ و دەستکاریکردن دەستپێبکەیت."
        },
        {
            "type": "hero_closing_links",
            "title": "داهاتووی ژیری دەستکرد & پەیوەندی",
            "subtitle": "سوپاس بۆ بەشداریتان لە کۆرسی ژیری دەستکرد (AI لە ژیانی ڕۆژانە و کاردا)",
            "presenter": "ئەندازیار عبدالرحمن اسماعیل",
            "links": [
                {"label": "تێلیگرام (شەخسی):", "value": "Aghay Andazyar (@Agha_ACE)"},
                {"label": "چەناڵی تێلیگرام:", "value": "داهاتووی تەکنەلۆجیا"},
                {"label": "وەتس ئاپ:", "value": "07504342452 (+9647504342452)"}
            ],
            "notes": "بەشی کۆتایی: سوپاسگوزاری، تێبینی کۆتایی، و پێدانی لینکەکانی پەیوەندی تێلیگرام، چەناڵ و وەتس ئاپ."
        }
    ]

    for idx, s_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Background fill
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

        # AI Artwork background overlay
        bg_art_map = {
            "hero": r"c:\Users\kurdn\Desktop\Semenar ChatGPT\assets\hero_ai_future.png",
            "hero_break": r"c:\Users\kurdn\Desktop\Semenar ChatGPT\assets\hero_ai_future.png",
            "hero_closing_links": r"c:\Users\kurdn\Desktop\Semenar ChatGPT\assets\hero_ai_future.png",
            "workflow_diag": r"c:\Users\kurdn\Desktop\Semenar ChatGPT\assets\ai_workflow_nodes.png",
            "grid5": r"c:\Users\kurdn\Desktop\Semenar ChatGPT\assets\ai_doc_analytics.png",
            "grid6_agenda": r"c:\Users\kurdn\Desktop\Semenar ChatGPT\assets\ai_doc_analytics.png",
            "comparison": r"c:\Users\kurdn\Desktop\Semenar ChatGPT\assets\hero_ai_future.png"
        }
        img_file = bg_art_map.get(s_data["type"])
        if img_file and os.path.exists(img_file):
            try:
                bg_pic = slide.shapes.add_picture(img_file, 0, 0, prs.slide_width, prs.slide_height)
                slide.shapes._spTree.remove(bg_pic._element)
                slide.shapes._spTree.insert(2, bg_pic._element)
            except Exception:
                pass

        # Speaker notes
        if s_data.get("notes"):
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.text = s_data["notes"]

        # Center Quote in Header across all slides
        quote_box = slide.shapes.add_textbox(Inches(4.0), Inches(0.35), Inches(5.33), Inches(0.4))
        tf_q = quote_box.text_frame
        tf_q.word_wrap = True
        p_q = tf_q.paragraphs[0]
        p_q.text = "« دەبێت کۆمپیوتەر بتوانێت بیربکاتەوە؟ »"
        p_q.font.name = FONT_NAME
        p_q.font.size = Pt(12)
        p_q.font.bold = True
        p_q.font.color.rgb = COLOR_CYAN
        p_q.alignment = PP_ALIGN.CENTER

        # Category Header
        if "category" in s_data:
            cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
            tf_cat = cat_box.text_frame
            tf_cat.word_wrap = True
            p_cat = tf_cat.paragraphs[0]
            p_cat.text = s_data["category"]
            p_cat.font.name = FONT_NAME
            p_cat.font.size = Pt(13)
            p_cat.font.bold = True
            p_cat.font.color.rgb = COLOR_CYAN
            p_cat.alignment = PP_ALIGN.RIGHT

        # Slide Title
        if "title" in s_data and s_data["type"] not in ["hero", "hero_break", "hero_closing_links"]:
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
            tf_t = title_box.text_frame
            tf_t.word_wrap = True
            p_t = tf_t.paragraphs[0]
            p_t.text = s_data["title"]
            p_t.font.name = FONT_NAME
            p_t.font.size = Pt(26)
            p_t.font.bold = True
            p_t.font.color.rgb = COLOR_WHITE
            p_t.alignment = PP_ALIGN.RIGHT

        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4))
        tf_f = footer_box.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = f"کۆرسی ژیری دەستکرد  •  سلاید {idx+1} لە {len(slides_data)}"
        p_f.font.name = FONT_NAME
        p_f.font.size = Pt(10)
        p_f.font.color.rgb = COLOR_MUTED
        p_f.alignment = PP_ALIGN.LEFT

        stype = s_data["type"]

        if stype in ["hero", "hero_break"]:
            hero_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.33), Inches(5.0))
            hero_card.fill.solid()
            hero_card.fill.fore_color.rgb = COLOR_CARD
            hero_card.line.color.rgb = COLOR_BLUE
            hero_card.line.width = Pt(1.5)

            tb = slide.shapes.add_textbox(Inches(1.8), Inches(1.8), Inches(9.73), Inches(3.8))
            tf = tb.text_frame
            tf.word_wrap = True

            if s_data.get("tag"):
                p0 = tf.paragraphs[0]
                p0.text = s_data["tag"]
                p0.font.name = FONT_NAME
                p0.font.size = Pt(14)
                p0.font.color.rgb = COLOR_CYAN
                p0.font.bold = True
                p0.alignment = PP_ALIGN.CENTER

            p1 = tf.add_paragraph() if s_data.get("tag") else tf.paragraphs[0]
            p1.text = s_data["title"]
            p1.font.name = FONT_NAME
            p1.font.size = Pt(36)
            p1.font.bold = True
            p1.font.color.rgb = COLOR_WHITE
            p1.alignment = PP_ALIGN.CENTER
            p1.space_before = Pt(14)

            p2 = tf.add_paragraph()
            p2.text = s_data["subtitle"]
            p2.font.name = FONT_NAME
            p2.font.size = Pt(18)
            p2.font.color.rgb = COLOR_LIGHT
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(14)

            if s_data.get("presenter"):
                p3 = tf.add_paragraph()
                p3.text = s_data["presenter"]
                p3.font.name = FONT_NAME
                p3.font.size = Pt(16)
                p3.font.bold = True
                p3.font.color.rgb = COLOR_CYAN
                p3.alignment = PP_ALIGN.CENTER
                p3.space_before = Pt(20)

        elif stype == "hero_closing_links":
            hero_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.0), Inches(10.93), Inches(5.5))
            hero_card.fill.solid()
            hero_card.fill.fore_color.rgb = COLOR_CARD
            hero_card.line.color.rgb = COLOR_CYAN
            hero_card.line.width = Pt(1.5)

            tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(10.33), Inches(5.1))
            tf = tb.text_frame
            tf.word_wrap = True

            p1 = tf.paragraphs[0]
            p1.text = s_data["title"]
            p1.font.name = FONT_NAME
            p1.font.size = Pt(32)
            p1.font.bold = True
            p1.font.color.rgb = COLOR_WHITE
            p1.alignment = PP_ALIGN.CENTER

            p2 = tf.add_paragraph()
            p2.text = s_data["subtitle"]
            p2.font.name = FONT_NAME
            p2.font.size = Pt(16)
            p2.font.color.rgb = COLOR_LIGHT
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(8)

            p3 = tf.add_paragraph()
            p3.text = s_data["presenter"]
            p3.font.name = FONT_NAME
            p3.font.size = Pt(16)
            p3.font.bold = True
            p3.font.color.rgb = COLOR_CYAN
            p3.alignment = PP_ALIGN.CENTER
            p3.space_before = Pt(10)

            # Add Links Info
            for l_info in s_data["links"]:
                pl = tf.add_paragraph()
                pl.text = f"🔗 {l_info['label']} {l_info['value']}"
                pl.font.name = FONT_NAME
                pl.font.size = Pt(15)
                pl.font.color.rgb = COLOR_WHITE
                pl.alignment = PP_ALIGN.CENTER
                pl.space_before = Pt(8)

        elif stype == "grid6_agenda":
            items = s_data["items"]
            card_w = Inches(3.6)
            card_h = Inches(2.1)
            gap_x = Inches(0.4)
            gap_y = Inches(0.3)
            start_x = Inches(0.8)
            start_y = Inches(1.8)

            for i, it in enumerate(items):
                col = i % 3
                row = i // 3
                cx = start_x + col * (card_w + gap_x)
                cy = start_y + row * (card_h + gap_y)

                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, card_w, card_h)
                card.fill.solid()
                card.fill.fore_color.rgb = COLOR_CARD
                card.line.color.rgb = COLOR_BLUE if i % 2 == 0 else COLOR_PURPLE
                card.line.width = Pt(1)

                tb = slide.shapes.add_textbox(cx + Inches(0.15), cy + Inches(0.15), card_w - Inches(0.3), card_h - Inches(0.3))
                tf = tb.text_frame
                tf.word_wrap = True

                p0 = tf.paragraphs[0]
                p0.text = it["num"]
                p0.font.name = FONT_NAME
                p0.font.size = Pt(16)
                p0.font.bold = True
                p0.font.color.rgb = COLOR_CYAN
                p0.alignment = PP_ALIGN.RIGHT

                p1 = tf.add_paragraph()
                p1.text = it["title"]
                p1.font.name = FONT_NAME
                p1.font.size = Pt(15)
                p1.font.bold = True
                p1.font.color.rgb = COLOR_WHITE
                p1.alignment = PP_ALIGN.RIGHT
                p1.space_before = Pt(4)

                p2 = tf.add_paragraph()
                p2.text = it["body"]
                p2.font.name = FONT_NAME
                p2.font.size = Pt(12)
                p2.font.color.rgb = COLOR_LIGHT
                p2.alignment = PP_ALIGN.RIGHT
                p2.space_before = Pt(4)

        elif stype == "grid5":
            items = s_data["items"]
            num_items = len(items)
            gap = Inches(0.3)
            card_w = (Inches(11.73) - (num_items - 1) * gap) / num_items
            start_x = Inches(0.8)

            for i, it in enumerate(items):
                cx = start_x + i * (card_w + gap)
                cy = Inches(1.8)
                ch = Inches(4.5)

                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, card_w, ch)
                card.fill.solid()
                card.fill.fore_color.rgb = COLOR_CARD
                card.line.color.rgb = COLOR_BLUE if i % 2 == 0 else COLOR_PURPLE
                card.line.width = Pt(1)

                tb = slide.shapes.add_textbox(cx + Inches(0.15), cy + Inches(0.2), card_w - Inches(0.3), ch - Inches(0.4))
                tf = tb.text_frame
                tf.word_wrap = True

                p1 = tf.paragraphs[0]
                p1.text = it["title"]
                p1.font.name = FONT_NAME
                p1.font.size = Pt(16)
                p1.font.bold = True
                p1.font.color.rgb = COLOR_CYAN if i % 2 == 0 else COLOR_WHITE
                p1.alignment = PP_ALIGN.RIGHT

                p2 = tf.add_paragraph()
                p2.text = it["body"]
                p2.font.name = FONT_NAME
                p2.font.size = Pt(13)
                p2.font.color.rgb = COLOR_LIGHT
                p2.alignment = PP_ALIGN.RIGHT
                p2.space_before = Pt(12)

        elif stype == "comparison":
            w_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5))
            w_card.fill.solid()
            w_card.fill.fore_color.rgb = RGBColor(0x2A, 0x15, 0x18)
            w_card.line.color.rgb = COLOR_DANGER
            w_card.line.width = Pt(1.5)

            tb_w = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.1))
            tf_w = tb_w.text_frame
            tf_w.word_wrap = True
            
            pw0 = tf_w.paragraphs[0]
            pw0.text = s_data["weak_title"]
            pw0.font.name = FONT_NAME
            pw0.font.size = Pt(18)
            pw0.font.bold = True
            pw0.font.color.rgb = COLOR_DANGER
            pw0.alignment = PP_ALIGN.RIGHT

            pw1 = tf_w.add_paragraph()
            pw1.text = s_data["weak_code"]
            pw1.font.name = FONT_NAME
            pw1.font.size = Pt(16)
            pw1.font.bold = True
            pw1.font.color.rgb = COLOR_WHITE
            pw1.alignment = PP_ALIGN.RIGHT
            pw1.space_before = Pt(14)

            pw2 = tf_w.add_paragraph()
            pw2.text = s_data["weak_body"]
            pw2.font.name = FONT_NAME
            pw2.font.size = Pt(14)
            pw2.font.color.rgb = COLOR_LIGHT
            pw2.alignment = PP_ALIGN.RIGHT
            pw2.space_before = Pt(14)

            s_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
            s_card.fill.solid()
            s_card.fill.fore_color.rgb = RGBColor(0x10, 0x2A, 0x20)
            s_card.line.color.rgb = COLOR_SUCCESS
            s_card.line.width = Pt(1.5)

            tb_s = slide.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.1))
            tf_s = tb_s.text_frame
            tf_s.word_wrap = True

            ps0 = tf_s.paragraphs[0]
            ps0.text = s_data["strong_title"]
            ps0.font.name = FONT_NAME
            ps0.font.size = Pt(18)
            ps0.font.bold = True
            ps0.font.color.rgb = COLOR_SUCCESS
            ps0.alignment = PP_ALIGN.RIGHT

            ps1 = tf_s.add_paragraph()
            ps1.text = s_data["strong_code"]
            ps1.font.name = FONT_NAME
            ps1.font.size = Pt(14)
            ps1.font.bold = True
            ps1.font.color.rgb = COLOR_WHITE
            ps1.alignment = PP_ALIGN.RIGHT
            ps1.space_before = Pt(10)

            ps2 = tf_s.add_paragraph()
            ps2.text = s_data["strong_body"]
            ps2.font.name = FONT_NAME
            ps2.font.size = Pt(14)
            ps2.font.color.rgb = COLOR_LIGHT
            ps2.alignment = PP_ALIGN.RIGHT
            ps2.space_before = Pt(10)

        elif stype == "presenter_intro":
            # Right Column Card (Profile Card)
            profile_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(1.8), Inches(4.2), Inches(4.5))
            profile_card.fill.solid()
            profile_card.fill.fore_color.rgb = COLOR_CARD
            profile_card.line.color.rgb = COLOR_BLUE
            profile_card.line.width = Pt(1.5)

            # Profile Picture inside Right Column
            img_path = s_data["image_path"]
            if os.path.exists(img_path):
                # Add picture
                slide.shapes.add_picture(img_path, Inches(9.3), Inches(2.2), Inches(2.2), Inches(2.2))

            # Name and details text inside Profile Card
            tb_p = slide.shapes.add_textbox(Inches(8.4), Inches(4.6), Inches(4.0), Inches(1.5))
            tf_p = tb_p.text_frame
            tf_p.word_wrap = True
            
            p_name = tf_p.paragraphs[0]
            p_name.text = s_data["name"]
            p_name.font.name = FONT_NAME
            p_name.font.size = Pt(20)
            p_name.font.bold = True
            p_name.font.color.rgb = COLOR_WHITE
            p_name.alignment = PP_ALIGN.CENTER
            
            p_role = tf_p.add_paragraph()
            p_role.text = s_data["roles"]
            p_role.font.name = FONT_NAME
            p_role.font.size = Pt(13)
            p_role.font.color.rgb = COLOR_CYAN
            p_role.alignment = PP_ALIGN.CENTER
            p_role.space_before = Pt(8)

            # Left Column (Education & Experience Cards)
            # Education Card (Card 1)
            ed_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(7.2), Inches(2.1))
            ed_card.fill.solid()
            ed_card.fill.fore_color.rgb = COLOR_CARD
            ed_card.line.color.rgb = COLOR_PURPLE
            ed_card.line.width = Pt(1)

            tb_ed = slide.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(6.8), Inches(1.9))
            tf_ed = tb_ed.text_frame
            tf_ed.word_wrap = True
            
            p_ed_title = tf_ed.paragraphs[0]
            p_ed_title.text = "🎓 بڕوانامەکان (Education)"
            p_ed_title.font.name = FONT_NAME
            p_ed_title.font.size = Pt(16)
            p_ed_title.font.bold = True
            p_ed_title.font.color.rgb = COLOR_PURPLE
            p_ed_title.alignment = PP_ALIGN.RIGHT
            
            for ed_item in s_data["education"]:
                p_item = tf_ed.add_paragraph()
                p_item.text = f"✦ {ed_item}"
                p_item.font.name = FONT_NAME
                p_item.font.size = Pt(13)
                p_item.font.color.rgb = COLOR_LIGHT
                p_item.alignment = PP_ALIGN.RIGHT
                p_item.space_before = Pt(4)

            # Experience Card (Card 2)
            exp_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.2), Inches(7.2), Inches(2.1))
            exp_card.fill.solid()
            exp_card.fill.fore_color.rgb = COLOR_CARD
            exp_card.line.color.rgb = COLOR_CYAN
            exp_card.line.width = Pt(1)

            tb_exp = slide.shapes.add_textbox(Inches(1.0), Inches(4.3), Inches(6.8), Inches(1.9))
            tf_exp = tb_exp.text_frame
            tf_exp.word_wrap = True
            
            p_exp_title = tf_exp.paragraphs[0]
            p_exp_title.text = "💼 ئەزموونی کارکردن (Experience)"
            p_exp_title.font.name = FONT_NAME
            p_exp_title.font.size = Pt(16)
            p_exp_title.font.bold = True
            p_exp_title.font.color.rgb = COLOR_CYAN
            p_exp_title.alignment = PP_ALIGN.RIGHT
            
            for exp_item in s_data["experience"]:
                p_item = tf_exp.add_paragraph()
                p_item.text = f"✦ {exp_item}"
                p_item.font.name = FONT_NAME
                p_item.font.size = Pt(13)
                p_item.font.color.rgb = COLOR_LIGHT
                p_item.alignment = PP_ALIGN.RIGHT
                p_item.space_before = Pt(4)

        elif stype == "workflow_diag":
            nodes = s_data["nodes"]
            nw = Inches(2.7)
            gap = Inches(0.3)
            start_x = Inches(0.8)

            for i, nd in enumerate(nodes):
                cx = start_x + i * (nw + gap)
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(2.2), nw, Inches(3.8))
                card.fill.solid()
                card.fill.fore_color.rgb = COLOR_CARD
                card.line.color.rgb = COLOR_SUCCESS if i == 3 else (COLOR_CYAN if i == 2 else COLOR_BLUE)
                card.line.width = Pt(1.5)

                tb = slide.shapes.add_textbox(cx + Inches(0.15), Inches(2.4), nw - Inches(0.3), Inches(3.4))
                tf = tb.text_frame
                tf.word_wrap = True

                p0 = tf.paragraphs[0]
                p0.text = f"هەنگاوی {nd['step']}"
                p0.font.name = FONT_NAME
                p0.font.size = Pt(14)
                p0.font.bold = True
                p0.font.color.rgb = COLOR_CYAN
                p0.alignment = PP_ALIGN.CENTER

                p1 = tf.add_paragraph()
                p1.text = nd["title"]
                p1.font.name = FONT_NAME
                p1.font.size = Pt(18)
                p1.font.bold = True
                p1.font.color.rgb = COLOR_WHITE
                p1.alignment = PP_ALIGN.CENTER
                p1.space_before = Pt(10)

                p2 = tf.add_paragraph()
                p2.text = nd["body"]
                p2.font.name = FONT_NAME
                p2.font.size = Pt(14)
                p2.font.color.rgb = COLOR_LIGHT
                p2.alignment = PP_ALIGN.CENTER
                p2.space_before = Pt(12)

    output_path = r"c:\Users\kurdn\Desktop\Semenar ChatGPT\سیمینارە_کۆتایی_Redesigned.pptx"
    prs.save(output_path)
    print(f"Presentation successfully created and saved at: {output_path}")

if __name__ == "__main__":
    import sys
    sys.stdout.reconfigure(encoding='utf-8')
    create_presentation()
